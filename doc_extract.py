"""Extract content from documents and build Gemini multimodal parts.

PDF  → uploaded to Gemini Files API (native page-as-image processing)
PPTX → slide text + embedded images sent inline
DOCX → plain text
TXT  → plain text
"""

from __future__ import annotations

import io
import os
import tempfile
import time
from pathlib import Path

from dotenv import load_dotenv
from google import genai
from google.genai import types

load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")


def _client() -> genai.Client:
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY manquante dans .env")
    return genai.Client(api_key=GEMINI_API_KEY)


def _extract_pptx(file_bytes: bytes, filename: str) -> list[types.Part]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[types.Part] = []
    slide_texts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    mime = shape.image.content_type or "image/png"
                    parts.append(types.Part.from_bytes(data=img_bytes, mime_type=mime))
                except Exception:
                    pass

        if texts:
            slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))

    if slide_texts:
        header = f"Document : {filename}\n\n" + "\n\n".join(slide_texts)
        parts.insert(0, types.Part.from_text(header))

    return parts


def _extract_docx(file_bytes: bytes, filename: str) -> list[types.Part]:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if not text:
        return []
    return [types.Part.from_text(f"Document : {filename}\n\n{text}")]


def _extract_txt(file_bytes: bytes, filename: str) -> list[types.Part]:
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")
    if not text.strip():
        return []
    return [types.Part.from_text(f"Document : {filename}\n\n{text}")]


def _upload_pdf(file_bytes: bytes, filename: str) -> list[types.Part]:
    client = _client()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        uploaded = client.files.upload(
            file=tmp_path,
            config=types.UploadFileConfig(mime_type="application/pdf", display_name=filename),
        )
        while uploaded.state.name == "PROCESSING":
            time.sleep(2)
            uploaded = client.files.get(name=uploaded.name)
        if uploaded.state.name == "FAILED":
            raise RuntimeError(f"Échec upload PDF Gemini : {filename}")
        return [
            types.Part.from_text(f"Document PDF : {filename}"),
            types.Part.from_uri(file_uri=uploaded.uri, mime_type="application/pdf"),
        ]
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def build_doc_parts(file_bytes: bytes, filename: str) -> list[types.Part]:
    """Return a list of Gemini Parts for the given document."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _upload_pdf(file_bytes, filename)
    if ext == ".pptx":
        return _extract_pptx(file_bytes, filename)
    if ext == ".docx":
        return _extract_docx(file_bytes, filename)
    if ext == ".txt":
        return _extract_txt(file_bytes, filename)
    return []
